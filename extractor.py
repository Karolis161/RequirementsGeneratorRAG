import fitz
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import docx


def extract_text_from_pdf(pdf_path):
    with fitz.open(pdf_path) as pdf:
        return "\n".join(page.get_text("text") for page in pdf if page.get_text("text")).strip()


def extract_text_from_ppt(ppt_path):
    text = []
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    return "\n".join(text).strip()


def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    return "\n".join(para.text.strip() for para in doc.paragraphs if para.text.strip())


def extract_text_from_website(url, timeout=10):
    try:
        response = requests.get(url, timeout=timeout)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        return "\n".join(p.get_text().strip() for p in soup.find_all("p") if p.get_text().strip())
    except requests.RequestException as e:
        print(f"⚠️ Error fetching {url}: {e}")
        return ""
